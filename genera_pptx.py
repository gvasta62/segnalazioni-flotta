# -*- coding: utf-8 -*-
# Crea una presentazione partendo dal template aziendale Busitalia
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import PP_PLACEHOLDER

TEMPLATE = r'D:/OneDrive - Gruppo Ferrovie Dello Stato/Busitalia_15032026/BusItalia_PPT_Template_03_26.pptx'
INFOG = 'infografica-progetto-orizzontale.png'
OUT = 'Presentazione_Segnalazione_Flotta.pptx'

prs = Presentation(TEMPLATE)

# --- rimuovi tutte le slide di esempio del template ---
sldIdLst = prs.slides._sldIdLst
for sld in list(sldIdLst):
    sldIdLst.remove(sld)

def layout(name):
    for l in prs.slide_layouts:
        if l.name == name:
            return l
    raise SystemExit('layout non trovato: ' + name)

def ph_by_type(slide, *types):
    for ph in slide.placeholders:
        if ph.placeholder_format.type in types:
            return ph
    return None

def set_title(slide, text):
    t = ph_by_type(slide, PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    if t is None and len(slide.placeholders):
        t = slide.placeholders[0]
    if t is not None:
        t.text = text
    return t

def first_body(slide, skip=()):
    for ph in slide.placeholders:
        if ph in skip:
            continue
        if ph.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.SUBTITLE,
                                          PP_PLACEHOLDER.OBJECT):
            return ph
    return None

# ===== Slide 1 — Copertina =====
s = prs.slides.add_slide(layout('1_Copertina'))
title = set_title(s, 'Segnalazione anomalie di flotta')
sub = first_body(s, skip=(title,))
if sub is not None:
    sub.text = 'Freccialink & The Mall — sistema digitale sviluppato internamente'

# ===== Slide 2 — Infografica (immagine a pagina intera) =====
s = prs.slides.add_slide(layout('Immagine pagina intera'))
pic_ph = ph_by_type(s, PP_PLACEHOLDER.PICTURE)
if pic_ph is not None:
    pic_ph.insert_picture(INFOG)
else:
    # fallback: immagine centrata
    s.shapes.add_picture(INFOG, Emu(0), Emu(int((prs.slide_height-prs.slide_width*1080/1920)/2)),
                         width=prs.slide_width)

# ===== Slide 3 — Punti chiave (titolo + bullet) =====
s = prs.slides.add_slide(layout('Titolo+Testo bulltes'))
set_title(s, 'Cosa rende misurabile il lavoro')
body = ph_by_type(s, PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
punti = [
    'Sviluppato internamente, a costo zero (0 €): nessun fornitore, licenza o canone.',
    'Carico di lavoro tracciato: volume delle segnalazioni per stato e per periodo.',
    'Tempi di risoluzione misurati: tempo medio dalla segnalazione alla chiusura.',
    'Reattività monitorata: rapidità con cui ogni segnalazione viene presa in carico.',
    'Notifiche email automatiche a ogni cambio di stato (apertura, presa in carico, chiusura).',
    'Storico completo e statistiche sempre aggiornate per la direzione.',
]
if body is not None:
    tf = body.text_frame
    tf.clear()
    for i, p in enumerate(punti):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = p
        para.level = 0

# ===== Slide 4 — Chiusura =====
s = prs.slides.add_slide(layout('Grazie - Fine'))
set_title(s, 'Grazie')

prs.save(OUT)
print('OK', OUT, '-', len(prs.slides.__iter__().__length_hint__() if hasattr(prs.slides,"__length_hint__") else 0))
print('slide:', len(prs.slides._sldIdLst))
